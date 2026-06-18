import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_presentation():
    # 1. Initialize Presentation
    prs = Presentation()

    # Define Brand Colors
    KENYA_RED = RGBColor(176, 0, 32)   # #B00020
    KENYA_GREEN = RGBColor(0, 102, 0)  # #006600
    KENYA_BLACK = RGBColor(0, 0, 0)

    # 2. Define Slide Content Data
    # This structure mirrors your "Hackathon_Presentation_Content.md"
    slides_data = [
        {
            "type": "title",
            "title": "KENYA DEBT GUARDIAN",
            "subtitle": "Empowering Fiscal Transparency & Sustainability with AI",
            "footer": "10Alytics Global Hackathon 2025 | [Your Team Name]",
            "visual_note": "Visual: High-quality background image of Nairobi skyline or Kenya map with digital nodes."
        },
        {
            "type": "content",
            "title": "THE PROBLEM - \"FLYING BLIND\"",
            "headline": "Fiscal Data is Fragmented & Opaque",
            "bullets": [
                "Data Silos: Critical financial data is buried in PDF reports, disconnected text files, and disparate excel sheets.",
                "Lack of Foresight: Historical reports tell us \"where we were,\" but fail to predict \"where we are going.\"",
                "Public Disconnect: Complex figures (Trillions) alienate the very citizens who bear the debt burden.",
                "The Result: Reactive policy-making instead of proactive strategy."
            ],
            "visual_note": "Visual: Icon of a scattered puzzle or a person looking confused at a stack of papers."
        },
        {
            "type": "content",
            "title": "THE SOLUTION",
            "headline": "A Centralized Fiscal Intelligence Hub",
            "bullets": [
                "Aggregated Truth: We combined data from Treasury, Central Bank, and KNBS into a single \"Master Dataset.\"",
                "Real-Time Analytics: Interactive dashboards that track Debt, Revenue, and Expenditure instantly.",
                "AI-Powered Forecasting: Using Facebook Prophet to predict future economic trajectories with 95% confidence intervals."
            ],
            "visual_note": "Visual: Screenshot of your Streamlit Dashboard (Home Page)."
        },
        {
            "type": "content",
            "title": "KEY INSIGHT 1 - THE DEBT BURDEN",
            "headline": "A Trajectory of Accelerated Borrowing",
            "bullets": [
                "Current Debt: [Insert Latest Value] Trillion KSh (Check Dashboard)",
                "Per Citizen Burden: ~[Insert Value] KSh for every Kenyan",
                "Trend: Annual debt accumulation has averaged [Insert Growth]% over the last 5 years.",
                "Insight: We are borrowing faster than we are growing."
            ],
            "visual_note": "Visual: The 'Historical Debt Trend Analysis' graph from your dashboard (Red Bars + Green Line)."
        },
        {
            "type": "content",
            "title": "KEY INSIGHT 2 - COMPOSITION RISKS",
            "headline": "The \"Double-Edged\" Debt Portfolio",
            "bullets": [
                "Domestic Debt (~53%): High interest rates crowd out the private sector, stifling local business growth.",
                "External Debt (~47%): Exposed to currency shocks. Every time the Shilling weakens, our debt stock automatically rises.",
                "Risk Level: HIGH. The near 50/50 split maximizes exposure to BOTH interest rate and exchange rate risks."
            ],
            "visual_note": "Visual: The 'Debt Composition Analysis' Pie Chart from your dashboard."
        },
        {
            "type": "content",
            "title": "KEY INSIGHT 3 - SUSTAINABILITY RED FLAG",
            "headline": "Operating in the \"Danger Zone\"",
            "bullets": [
                "Current Debt-to-GDP: [Insert Value]% (e.g., 68.4%)",
                "IMF Safe Limit: 55% (Green Line on Graph)",
                "Critical Threshold: 70% (Red Line on Graph)",
                "Status: We have breached the safe limit and are approaching the critical insolvency threshold."
            ],
            "visual_note": "Visual: The 'Debt Sustainability Analysis' Line Graph showing the breach of the 55% threshold."
        },
        {
            "type": "content",
            "title": "THE \"KILLER FEATURE\" - AI FORECAST",
            "headline": "Predicting the 2028 Fiscal Gap",
            "bullets": [
                "The Prediction: By 2028, Kenya is projected to accumulate an ADDITIONAL [Insert Forecast Value] Trillion KSh in debt.",
                "Revenue vs. Expenditure: Spending is projected to grow 1.5x faster than revenue collection.",
                "The Gap: A projected fiscal deficit of [Insert Value] Trillion KSh by 2028 if no action is taken."
            ],
            "visual_note": "Visual: The 'Projected Debt Trajectory' graph with the confidence interval shading."
        },
        {
            "type": "content",
            "title": "STRATEGIC RECOMMENDATIONS",
            "headline": "A Path to Sustainability (3-Pronged Strategy)",
            "bullets": [
                "Immediate (Stop the Bleeding): Freeze non-essential recurrent expenditure and negotiate debt service suspension.",
                "Medium Term (Heal the Patient): Expand tax base via digital compliance (not higher rates) and privatize non-strategic SOEs.",
                "Long Term (Build Immunity): Constitutional debt ceiling capped at 55% of GDP and export-led growth."
            ],
            "visual_note": "Visual: Icons representing each phase (Stopwatch, Chart, Shield)."
        },
        {
            "type": "content",
            "title": "IMPACT & SCALABILITY",
            "headline": "Beyond a Dashboard",
            "bullets": [
                "For Citizens: Translates \"Trillions\" into personal impact metrics (Debt per Person).",
                "For Investors: Provides an independent, data-driven risk assessment.",
                "For Policymakers: An early warning system to prevent fiscal crises before they happen.",
                "Scalability: Can be adapted for any African nation facing similar debt challenges (Ghana, Zambia, etc.)."
            ],
            "visual_note": "Visual: Map of Africa or icons of different user groups."
        },
        {
            "type": "content",
            "title": "CONCLUSION",
            "headline": "The Time for Action is Now",
            "bullets": [
                "Sustainable Debt = Prosperous Kenya.",
                "We cannot change the past borrowing, but with Kenya Debt Guardian, we can reshape our financial future."
            ],
            "visual_note": "Visual: Strong closing image (Kenya Flag or hopeful imagery).",
            "footer_override": "Thank You! | Q&A"
        }
    ]

    # 3. Iterate and Create Slides
    for i, data in enumerate(slides_data):
        if data['type'] == 'title':
            # Use Title Slide Layout (Index 0)
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set Title
            title = slide.shapes.title
            title.text = data['title']
            title.text_frame.paragraphs[0].font.color.rgb = KENYA_RED
            title.text_frame.paragraphs[0].font.bold = True
            
            # Set Subtitle
            subtitle = slide.placeholders[1]
            subtitle.text = data['subtitle'] + "\n\n" + data['footer']
            
            # Add Image if available (e.g., image_62ae41.jpg)
            img_path = 'image_62ae41.jpg'
            if os.path.exists(img_path):
                # Add image to the right side or background
                # Here we add it as a small logo/accent for style
                left = Inches(1)
                top = Inches(5)
                height = Inches(2)
                slide.shapes.add_picture(img_path, left, top, height=height)

        elif data['type'] == 'content':
            # Use Title and Content Layout (Index 1)
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set Title
            title = slide.shapes.title
            title.text = data['title']
            title.text_frame.paragraphs[0].font.color.rgb = KENYA_GREEN
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
            
            # Set Headline (as a text box below title)
            left = Inches(0.5)
            top = Inches(1.3)
            width = Inches(9)
            height = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = data['headline']
            p.font.color.rgb = KENYA_RED
            p.font.size = Pt(24)
            p.font.bold = True
            
            # Set Bullets (Main Body)
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear() # Clear default empty paragraph
            
            for bullet in data['bullets']:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
                p.space_after = Pt(14)
                
            # Add Placeholder Box for Visuals
            # This creates a grey box where you should paste your graph screenshots later
            left = Inches(6)
            top = Inches(2)
            width = Inches(3.5)
            height = Inches(3.5)
            shape = slide.shapes.add_shape(
                1, left, top, width, height # 1 is msoShapeRectangle
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            shape.line.color.rgb = RGBColor(200, 200, 200)
            shape.text = "[PASTE GRAPH HERE]\n" + data.get('visual_note', '')

        # Add Speaker Notes
        if 'visual_note' in data:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = f"VISUAL IDEA: {data['visual_note']}"

    # 4. Save Presentation
    output_filename = 'Kenya_Debt_Guardian_Pitch.pptx'
    prs.save(output_filename)
    print(f"Successfully generated {output_filename}!")

if __name__ == "__main__":
    create_presentation()